#!/usr/bin/env python

import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import time

SLD_LAYOUT_TITLE = 0
SLD_LAYOUT_TITLE_AND_CONTENT = 1
SLD_LAYOUT_SECTION = 2


# Create a new presentations
prs = Presentation()

# Add Header Slide
title_slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Developer Tools"
subtitle.text = "Tools to get you ready to do awesome things"


# Add Agenda Slide
slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
shapes = slide.shapes
shapes.title.text = "Agenda"
body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Markdown"
p = tf.add_paragraph()
p.text = "IDE"

#tf.text = [ "Markdown", "IDE", "iTerm", "Git", "Continous Integration"]

prs.save('Developer Tools Presentation.pptx')
